from pptx import Presentation
import json
import os
from pathlib import Path

# Load the presentation
prs = Presentation('DVM_Annual_Report-Final2.pptx')

achievers = []
output_dir = Path('extracted_achievers')
output_dir.mkdir(exist_ok=True)

# Iterate through slides
for slide_idx, slide in enumerate(prs.slides):
    slide_data = {
        'slide_number': slide_idx + 1,
        'text': [],
        'images': []
    }
    
    # Extract text
    for shape in slide.shapes:
        if hasattr(shape, "text"):
            if shape.text.strip():
                slide_data['text'].append(shape.text.strip())
        
        # Extract images
        if shape.shape_type == 13:  # Picture
            image = shape.image
            image_bytes = image.blob
            image_filename = f'slide_{slide_idx + 1}_img_{len(slide_data["images"]) + 1}.{image.ext}'
            image_path = output_dir / image_filename
            
            with open(image_path, 'wb') as f:
                f.write(image_bytes)
            
            slide_data['images'].append(image_filename)
    
    if slide_data['text'] or slide_data['images']:
        achievers.append(slide_data)

# Save to JSON
with open('extracted_achievers.json', 'w', encoding='utf-8') as f:
    json.dump(achievers, f, indent=2, ensure_ascii=False)

print(f"Extracted {len(achievers)} slides")
print("Data saved to extracted_achievers.json")
print(f"Images saved to {output_dir}")
